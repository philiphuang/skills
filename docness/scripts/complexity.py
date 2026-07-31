"""复杂度判定模块——提取文档复杂度指标，决定路由策略。

路由策略：
- simple: pandoc / pypdf / pandas+openpyxl（本地快速管道）
- complex: MinerU Skill（高保真转换）
"""

from __future__ import annotations

import json
import sys
from pathlib import Path


# 各格式的简单/复杂判定阈值
# 旧格式(-legacy) 在 converter.py 先经 LibreOffice 转为新格式，再用新格式阈值判定，
# 因此无需为 *-legacy 重复设阈值。
THRESHOLDS: dict[str, dict[str, int]] = {
    "pdf": {"pages": 50, "tables": 5, "images": 10},
    "word": {"pages_estimate": 80, "tables": 10, "images": 20},
    "excel": {"sheets": 3, "rows": 5000, "formulas": 100},
    "csv": {"rows": 10000, "columns": 20},
    "powerpoint": {"slides": 50, "images": 20},
}


def get_thresholds(file_type: str) -> dict[str, int]:
    """获取指定文件类型的复杂度阈值。

    Args:
        file_type: 文件类型标识（如 pdf/word/excel 等）

    Returns:
        阈值 dict，未定义的类型返回空 dict（不设限制，全走 simple）
    """
    return THRESHOLDS.get(file_type, {})


def extract_metrics_pdf(filepath: str | Path) -> dict:
    """提取 PDF 文件的复杂度指标。

    Returns:
        dict with keys: pages, tables (估算), images (估算), size_bytes
    """
    path = Path(filepath)
    metrics = {"pages": 0, "tables": 0, "images": 0, "size_bytes": path.stat().st_size}

    try:
        from pypdf import PdfReader

        reader = PdfReader(str(path))
        metrics["pages"] = len(reader.pages)
    except Exception:
        pass

    return metrics


def extract_metrics_docx(filepath: str | Path) -> dict:
    """提取 DOCX 文件的复杂度指标。

    Returns:
        dict with keys: pages_estimate, tables, images, size_bytes
    """
    path = Path(filepath)
    metrics = {
        "pages_estimate": 0,
        "tables": 0,
        "images": 0,
        "size_bytes": path.stat().st_size,
    }

    try:
        import docx

        doc = docx.Document(str(path))
        metrics["tables"] = len(doc.tables)
        metrics["images"] = sum(
            1 for r in doc.inline_shapes
            if r.type == docx.enum.shape.WD_INLINE_SHAPE_TYPE.PICTURE
        )
        # 粗略页数估算：每 3000 字符约 1 页
        char_count = sum(
            len(p.text) for para in doc.paragraphs for p in [para]
        )
        metrics["pages_estimate"] = max(1, char_count // 3000)
    except Exception:
        pass

    return metrics


def extract_metrics_pptx(filepath: str | Path) -> dict:
    """提取 PPTX 文件的复杂度指标。

    Returns:
        dict with keys: slides, images, size_bytes
    """
    path = Path(filepath)
    metrics = {"slides": 0, "images": 0, "size_bytes": path.stat().st_size}

    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        metrics["slides"] = len(prs.slides)
        img_count = 0
        for slide in prs.slides:
            img_count += sum(
                1 for shape in slide.shapes if shape.shape_type == 13  # MSO_SHAPE_TYPE.PICTURE
            )
        metrics["images"] = img_count
    except Exception:
        pass

    return metrics


def extract_metrics_xlsx(filepath: str | Path) -> dict:
    """提取 XLSX 文件的复杂度指标。

    Returns:
        dict with keys: sheets, rows, formulas, size_bytes
    """
    path = Path(filepath)
    metrics = {"sheets": 0, "rows": 0, "formulas": 0, "size_bytes": path.stat().st_size}

    try:
        import openpyxl

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=False)
        metrics["sheets"] = len(wb.sheetnames)
        total_rows = 0
        formulas = 0
        for sheet in wb:
            total_rows += sheet.max_row or 0
            # 计数公式（简单估算：含 = 的单元格）
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        formulas += 1
        wb.close()
        metrics["rows"] = total_rows
        metrics["formulas"] = formulas
    except Exception:
        pass

    return metrics


def extract_metrics(filepath: str | Path, file_type: str) -> dict:
    """根据文件类型提取复杂度指标。

    Args:
        filepath: 文件路径
        file_type: 文件类型标识

    Returns:
        复杂度指标 dict
    """
    extractors = {
        "pdf": extract_metrics_pdf,
        "word": extract_metrics_docx,
        "powerpoint": extract_metrics_pptx,
        "excel": extract_metrics_xlsx,
    }

    extractor = extractors.get(file_type)
    if extractor:
        return extractor(filepath)

    # 不支持的文件类型（或文件不存在），返回空指标
    path = Path(filepath)
    if path.exists():
        return {"size_bytes": path.stat().st_size}
    return {}


def is_complex(metrics: dict, file_type: str) -> tuple[bool, list[str]]:
    """判定文件是否属于复杂文档。

    Args:
        metrics: extract_metrics 返回的指标 dict
        file_type: 文件类型标识

    Returns:
        (is_complex, reasons): 是否复杂 + 超标项列表
    """
    thresholds = get_thresholds(file_type)
    if not thresholds:
        return False, []

    reasons: list[str] = []
    for key, threshold in thresholds.items():
        actual = metrics.get(key, 0)
        if actual > threshold:
            reasons.append(f"{key}={actual}(>{threshold})")

    return len(reasons) > 0, reasons


def decide_route(
    filepath: str | Path, file_type: str
) -> dict:
    """完整决策：提取指标 + 判定复杂度 + 输出路由决策。

    Returns:
        {
            "file_type": str,
            "metrics": dict,
            "is_complex": bool,
            "reasons": list[str],
            "route": "simple" | "complex",
            "recommended_tool": str,
        }
    """
    metrics = extract_metrics(filepath, file_type)
    complex_flag, reasons = is_complex(metrics, file_type)

    # 路由映射：简单/复杂各自用什么工具
    SIMPLE_TOOLS: dict[str, str] = {
        "pdf": "pandoc → pypdf",
        "word": "pandoc",
        "excel": "pandas+openpyxl",
        "csv": "pandas",
        "powerpoint": "pandoc",
    }
    COMPLEX_TOOLS: dict[str, str] = {
        "pdf": "MinerU",
        "word": "MinerU",
        "excel": "MinerU",
        "powerpoint": "MinerU",
    }

    if complex_flag:
        route = "complex"
        tool = COMPLEX_TOOLS.get(file_type, "MinerU")
    else:
        route = "simple"
        tool = SIMPLE_TOOLS.get(file_type, "pandoc")

    return {
        "file_type": file_type,
        "metrics": metrics,
        "is_complex": complex_flag,
        "reasons": reasons,
        "route": route,
        "recommended_tool": tool,
    }


def main():
    """命令行入口：python3 -m scripts.complexity <文件路径> <文件类型>。

    输出完整路由决策 JSON（metrics + is_complex + route + recommended_tool）。
    """
    if len(sys.argv) < 3:
        print(
            "用法: python3 -m scripts.complexity <文件路径> <文件类型>\n"
            "文件类型: pdf / word / excel / csv / powerpoint",
            file=sys.stderr,
        )
        sys.exit(2)

    filepath, file_type = sys.argv[1], sys.argv[2]
    decision = decide_route(filepath, file_type)
    print(json.dumps(decision, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
